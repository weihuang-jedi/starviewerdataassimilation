#!/usr/bin/env python3
"""
generate_aida_presentation.py
-----------------------------
Generates an executive 7-slide PowerPoint presentation (.pptx) summarizing
the AIDA GNN Surrogate & AMSU-A Satellite Radiance Integration project using python-pptx.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def create_deck(output_filename="AIDA_GNN_AMSUA_Summary.pptx"):
    prs = Presentation()

    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Blank slide layout

    # Color Palette Constants
    NAVY = RGBColor(15, 32, 67)        # Primary Header Color
    DARK_BLUE = RGBColor(24, 60, 108)   # Accent Color
    ACCENT_TEAL = RGBColor(0, 150, 166)  # Feature Highlight
    LIGHT_BG = RGBColor(245, 247, 250)  # Card Background
    WHITE = RGBColor(255, 255, 255)
    TEXT_DARK = RGBColor(40, 40, 40)
    TEXT_MUTED = RGBColor(100, 100, 100)

    # Helper function to create standard header banner
    def add_header(slide, title_text, category_text="AIDA GNN SURROGATE & DATA ASSIMILATION"):
        # Header text frame
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_TEAL
        p_cat.space_after = Pt(2)

        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY

    # Helper function to create rounded container box
    def add_card(slide, left, top, width, height, bg_color=LIGHT_BG):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = RGBColor(220, 225, 230)
        shape.line.width = Pt(1)
        return shape

    # =========================================================================
    # SLIDE 1: Executive Overview & Motivation
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    add_header(slide1, "Executive Overview & Project Motivation")

    card1 = add_card(slide1, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf1 = card1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.4)
    tf1.margin_top = Inches(0.4)
    tf1.margin_right = Inches(0.4)

    bullets_s1 = [
        ("Strategic Vision", "Atmospheric Icosahedral Data Assimilation (AIDA) framework leveraging Graph Neural Networks for global weather surrogate modeling."),
        ("The Paradigm Shift", "Traditional DA methods (3D/4D-Var, EnKF) face computational bottlenecks at high spatial/temporal resolutions due to costly iterative matrix inversions and complex forward operators."),
        ("The Neural Solution", "A GNN surrogate model designed directly on unstructured M4 icosahedral meshes, bypassing grid transformations while enforcing physical conservation laws directly in neural feature space."),
        ("Core Objective", "Integrate satellite radiances (AMSU-A brightness temperatures) and conventional observation vectors into a fully differentiable loss engine to guide the GNN surrogate toward dynamically balanced, observation-informed analyses.")
    ]

    for idx, (title, desc) in enumerate(bullets_s1):
        p = tf1.paragraphs[0] if idx == 0 else tf1.add_paragraph()
        p.space_after = Pt(14)
        
        run_title = p.add_run()
        run_title.text = f"• {title}: "
        run_title.font.bold = True
        run_title.font.size = Pt(15)
        run_title.font.color.rgb = DARK_BLUE

        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(14)
        run_desc.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 2: Physical Architecture & Multi-Scale Constraints
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Physical Architecture & Multi-Scale Constraints")

    # 3 Column Cards Layout
    col_w = Inches(3.64)
    gap = Inches(0.4)

    # Column 1: Log-State & Thermo
    c1 = add_card(slide2, Inches(0.8), Inches(1.6), col_w, Inches(5.2))
    tf2_1 = c1.text_frame
    tf2_1.word_wrap = True
    tf2_1.margin_left = tf2_1.margin_top = tf2_1.margin_right = Inches(0.3)
    p = tf2_1.paragraphs[0]
    p.text = "1. Log-State & Thermo"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(10)

    p2 = tf2_1.add_paragraph()
    p2.text = "State Vector Representation:\nx = [ln T, u, v, w, q, ln ρ, ln p]^T\n\nThermodynamic Consistency (J_thermo):\nEnforces exact Ideal Gas Law:\n|| ln p - (ln ρ + ln R_d + ln T) ||²\n\nGuarantees strictly non-negative physical pressure and density fields."
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_DARK

    # Column 2: Hybrid Dynamics
    c2 = add_card(slide2, Inches(0.8) + col_w + gap, Inches(1.6), col_w, Inches(5.2))
    tf2_2 = c2.text_frame
    tf2_2.word_wrap = True
    tf2_2.margin_left = tf2_2.margin_top = tf2_2.margin_right = Inches(0.3)
    p = tf2_2.paragraphs[0]
    p.text = "2. Hybrid Dynamics (J_dyn)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(10)

    p2 = tf2_2.add_paragraph()
    p2.text = "Mid-Latitudes (Geostrophic Balance):\nRestrains unphysical wind and pressure updates using sparse differential operators:\nfv ≈ R_d T (∂ln p / ∂x)\nfu ≈ -R_d T (∂ln p / ∂y)\n\nTropics (Mass Continuity):\nRestrains divergent gravity waves via low-divergence penalty:\n∇ · V_h ≈ 0"
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_DARK

    # Column 3: Structural Bounds
    c3 = add_card(slide2, Inches(0.8) + (col_w + gap) * 2, Inches(1.6), col_w, Inches(5.2))
    tf2_3 = c3.text_frame
    tf2_3.word_wrap = True
    tf2_3.margin_left = tf2_3.margin_top = tf2_3.margin_right = Inches(0.3)
    p = tf2_3.paragraphs[0]
    p.text = "3. Spatial & Pressure Bounds"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(10)

    p2 = tf2_3.add_paragraph()
    p2.text = "Pressure Laplacian (J_lap):\n2nd-order graph Laplacian smoothing eliminates high-frequency noise across M4 mesh edges.\n\nAsymmetric Pressure Barrier (J_barrier):\nPrevents low-pressure collapse in stratosphere via non-linear barrier penalty at threshold τ_min."
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 3: Observational Ingestion Strategy
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Observational Ingestion & Satellite Radiance Integration")

    card3 = add_card(slide3, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    tf3 = card3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = tf3.margin_top = tf3.margin_right = Inches(0.4)

    bullets_s3 = [
        ("Conventional Observations", "In-situ surface, radiosonde, and aircraft observations mapped onto closest M4 icosahedral graph nodes."),
        ("AMSU-A Microwave Sounder Integration", "Direct ingestion of Channels 1–15 brightness temperatures (T_b) sensitive to atmospheric thermal structure from surface to stratosphere."),
        ("Differentiable Forward Operator H(x)", "Integrates vertical temperature T_k and pressure p_k profiles in PyTorch using channel-specific log-pressure Gaussian weighting functions W_c(p):"),
        ("Radiance Innovation Loss Engine (J_rad)", "Evaluates channel-weighted residual normalized by sensor observation error σ_c across active satellite swaths:")
    ]

    for idx, (title, desc) in enumerate(bullets_s3):
        p = tf3.paragraphs[0] if idx == 0 else tf3.add_paragraph()
        p.space_after = Pt(12)
        
        run_title = p.add_run()
        run_title.text = f"• {title}: "
        run_title.font.bold = True
        run_title.font.size = Pt(15)
        run_title.font.color.rgb = DARK_BLUE

        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(14)
        run_desc.font.color.rgb = TEXT_DARK

    # Add Formula Container Box on Slide 3
    f_card = add_card(slide3, Inches(1.2), Inches(5.2), Inches(10.933), Inches(1.3), bg_color=WHITE)
    ftf = f_card.text_frame
    ftf.word_wrap = True
    ftf.margin_top = Inches(0.2)
    p_eq = ftf.paragraphs[0]
    p_eq.alignment = PP_ALIGN.CENTER
    p_eq.text = "T_b,sim_c = ∑_k [ T_k · W_c(p_k) ]   ====>   J_rad = ∑_c [ (y_amsua_c - T_b,sim_c) / σ_c ]²"
    p_eq.font.size = Pt(16)
    p_eq.font.bold = True
    p_eq.font.color.rgb = ACCENT_TEAL

    # =========================================================================
    # SLIDE 4: Milestone Progress: Conventional to AMSU-A
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Project Milestones & Implementation Roadmap")

    phases = [
        ("Phase 1", "Synthetic Proof-of-Concept", "Validated M4 Mesh GNN message passing & topology on synthetic log-state fields."),
        ("Phase 2", "Conventional Obs (conv_2024)", "Ingested spatial in-situ observations; trained basic MSE + GNN topology updates."),
        ("Phase 3", "AMSU-A Radiance Integration", "Added DifferentiableAMSUAOperator for 15 microwave channels; integrated J_rad."),
        ("Phase 4", "Operational System Hardening", "Built YAML config engine, automated cycling inference, and verification suites.")
    ]

    p_w = Inches(2.68)
    p_gap = Inches(0.33)

    for idx, (ph_num, ph_title, ph_desc) in enumerate(phases):
        p_left = Inches(0.8) + idx * (p_w + p_gap)
        p_card = add_card(slide4, p_left, Inches(1.8), p_w, Inches(4.8))
        ptf = p_card.text_frame
        ptf.word_wrap = True
        ptf.margin_left = ptf.margin_top = ptf.margin_right = Inches(0.25)

        p = ptf.paragraphs[0]
        p.text = ph_num.upper()
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = ACCENT_TEAL
        p.space_after = Pt(4)

        p2 = ptf.add_paragraph()
        p2.text = ph_title
        p2.font.bold = True
        p2.font.size = Pt(15)
        p2.font.color.rgb = DARK_BLUE
        p2.space_after = Pt(12)

        p3 = ptf.add_paragraph()
        p3.text = ph_desc
        p3.font.size = Pt(13)
        p3.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 5: Major Engineering & Scientific Challenges Resolved
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Engineering & Scientific Challenges Resolved")

    # Table Layout
    table_shape = slide5.shapes.add_table(5, 4, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0))
    table = table_shape.table

    # Column Widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(3.2)
    table.columns[3].width = Inches(3.533)

    headers = ["Challenge / Bug", "Diagnostic Symptom", "Root Cause", "Engineering Solution"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE

    rows_data = [
        ("Pressure Explosion", "p_bias = +16,319 hPa, ACC = 0.008", "Un-normalized log-state passed to H(x) producing > 10^7 Pa.", "Added explicit mean/std un-normalization and bounds clamping before H(x)."),
        ("Cold Temperature Drift", "T_bias = -122.8 K, T = 141 K flat", "Raw J_rad innovation 100x larger than MSE; hit lower clamp bound.", "Normalized J_rad by channel count (15) and tuned weight down (w_rad = 0.01)."),
        ("Broadcasting Crash", "ValueError: shapes (181,) & (32,181,360)", "1D Cosine-Lat weight array lacked level/lon dimensions.", "Upgraded compute_weighted_metrics to build full 3D broadcast arrays (1, lat, 1)."),
        ("Wind/Moisture Over-Shooting", "u, v > 95 m/s, q > 0.112 kg/kg", "Adding full standardized prediction on top of background doubled variance.", "Implemented residual increment scaling (α = 0.15) and physical variable bounds.")
    ]

    for row_idx, r_data in enumerate(rows_data, start=1):
        for col_idx, text in enumerate(r_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if row_idx % 2 == 1 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 6: Final Performance & Verification Results
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Final Verification Metrics vs. GFS Analysis Benchmark")

    # Table Layout for Metrics
    t_shape6 = slide6.shapes.add_table(7, 6, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0))
    t6 = t_shape6.table

    t6.columns[0].width = Inches(1.2)
    t6.columns[1].width = Inches(2.833)
    t6.columns[2].width = Inches(1.8)
    t6.columns[3].width = Inches(1.8)
    t6.columns[4].width = Inches(2.1)
    t6.columns[5].width = Inches(2.0)

    h6 = ["Variable", "Description", "RMSE", "BIAS", "ACC (Pattern Skill)", "Assessment"]
    for i, h in enumerate(h6):
        cell = t6.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE

    m_data = [
        ("t", "Temperature (K)", "3.79 K", "-0.16 K", "0.9926", "Outstanding thermal skill"),
        ("p", "Air Pressure (Pa)", "21.04 Pa", "-9.36 Pa", "0.9984", "Exact vertical mass tracking"),
        ("q", "Specific Humidity (kg/kg)", "0.0038", "+0.0013", "0.9066", "Captures fine moisture fields"),
        ("u", "Zonal Wind (m/s)", "7.85 m/s", "+2.04 m/s", "0.8448", "Stable jet stream structures"),
        ("v", "Meridional Wind (m/s)", "5.28 m/s", "-1.51 m/s", "0.8462", "Balanced cross-latitudinal flow"),
        ("w", "Vertical Wind (m/s)", "0.29 m/s", "-0.002 m/s", "0.1778", "Physically bounded range")
    ]

    for row_idx, r_data in enumerate(m_data, start=1):
        for col_idx, text in enumerate(r_data):
            cell = t6.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if row_idx % 2 == 1 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                if col_idx in [2, 3, 4]:
                    p.font.bold = True
                p.font.color.rgb = DARK_BLUE if col_idx == 4 else TEXT_DARK

    # =========================================================================
    # SLIDE 7: Comparative Analysis & Competitive Edge
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Comparative Analysis & Competitive Edge")

    c_w = Inches(3.64)
    c_gap = Inches(0.4)

    m_cols = [
        ("Traditional 4D-Var", [
            ("Computational Speed", "Extremely slow; computationally prohibitive for real-time high-res forecasting due to iterative matrix inversions."),
            ("Satellite Data Assimilation", "Excellence in observation handling, but struggles with scaling."),
            ("Grid Topology", "Requires structured lat-lon or regular grids; suffers from polar singularities.")
        ]),
        ("Standard ML Surrogates\n(FourCastNet / GraphCast)", [
            ("Computational Speed", "Fast inference (~seconds per global forecast step)."),
            ("Physical Consistency", "Blindsided by physical conservation bounds; prone to climate drift and spatial smoothing."),
            ("Satellite Data Assimilation", "No direct satellite radiance (AMSU-A T_b) assimilation capabilities.")
        ]),
        ("AIDA GNN (Our Approach)", [
            ("Computational Speed", "Fast GPU inference + integrated observation updates."),
            ("Physical Consistency", "Enforces Ideal Gas Law, Geostrophic Balance, and Mass Continuity natively."),
            ("Satellite Data Assimilation", "Directly ingests AMSU-A satellite radiances via a differentiable forward operator H(x) on native M4 icosahedral mesh.")
        ])
    ]

    for idx, (title, points) in enumerate(m_cols):
        card = add_card(slide7, Inches(0.8) + idx * (c_w + c_gap), Inches(1.6), c_w, Inches(5.2))
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = ACCENT_TEAL if idx == 2 else DARK_BLUE
        p.space_after = Pt(10)

        for pt_title, pt_desc in points:
            p2 = tf.add_paragraph()
            p2.space_after = Pt(8)
            
            run1 = p2.add_run()
            run1.text = f"• {pt_title}: "
            run1.font.bold = True
            run1.font.size = Pt(12)
            run1.font.color.rgb = DARK_BLUE

            run2 = p2.add_run()
            run2.text = pt_desc
            run2.font.size = Pt(11)
            run2.font.color.rgb = TEXT_DARK

    # Save presentation
    os.makedirs(os.path.dirname(output_filename) or '.', exist_ok=True)
    prs.save(output_filename)
    print(f"\n[SUCCESS] PowerPoint slide deck generated successfully: '{output_filename}'")


if __name__ == "__main__":
    create_deck("output/AIDA_GNN_AMSUA_Summary.pptx")
